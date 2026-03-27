"""
Agent 2 - Knowledge & Compliance (RAG + Vector DB)

Professional but simple design:
- Persistent Chroma vector database on disk
- Knowledge sync from folder (supports .txt, .md, .pdf, .ppt, .pptx, .xlsx)
- OCR for PDFs using Tesseract
- Grounded answer generation with Ollama
- Basic compliance checks and quality metrics
"""

from __future__ import annotations

import argparse
import json
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List

import chromadb
from ollama import Client

try:
    import pytesseract
    from pdf2image import convert_from_path
    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from openpyxl import load_workbook
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False


SUPPORTED_EXTENSIONS = {".txt", ".md", ".pdf", ".ppt", ".pptx", ".xlsx"}
DEFAULT_KB_DIR = Path("backend/data/rag_docs")
DEFAULT_DB_DIR = Path("backend/data/vector_db")
DEFAULT_COLLECTION = "product_knowledge"


def _chunk_text(text: str, chunk_size: int = 900, overlap: int = 150) -> List[str]:
    clean = re.sub(r"\s+", " ", text).strip()
    if not clean:
        return []

    chunks: List[str] = []
    start = 0
    n = len(clean)
    while start < n:
        end = min(start + chunk_size, n)
        chunks.append(clean[start:end])
        if end >= n:
            break
        start = max(0, end - overlap)
    return chunks


def _extract_text_from_pdf(file_path: Path) -> str:
    """Extract text from PDF using OCR (Tesseract)."""
    if not TESSERACT_AVAILABLE:
        return f"[PDF file skipped - pytesseract/pdf2image not available: {file_path.name}]"
    try:
        images = convert_from_path(str(file_path))
        text_parts = []
        for img in images:
            text = pytesseract.image_to_string(img)
            text_parts.append(text)
        return "\n".join(text_parts)
    except Exception as e:
        return f"[Error extracting PDF {file_path.name}: {str(e)}]"


def _extract_text_from_pptx(file_path: Path) -> str:
    """Extract text from PowerPoint (.pptx) file."""
    if not PPTX_AVAILABLE:
        return f"[PPT file skipped - python-pptx not available: {file_path.name}]"
    try:
        presentation = Presentation(str(file_path))
        text_parts = []
        for slide_num, slide in enumerate(presentation.slides, 1):
            text_parts.append(f"--- Slide {slide_num} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
        return "\n".join(text_parts)
    except Exception as e:
        return f"[Error extracting PPT {file_path.name}: {str(e)}]"


def _extract_text_from_ppt(file_path: Path) -> str:
    """Extract text from PowerPoint (.ppt) file - delegate to .pptx handler."""
    # .ppt and .pptx use the same library
    return _extract_text_from_pptx(file_path)


def _extract_text_from_xlsx(file_path: Path) -> str:
    """Extract text from Excel (.xlsx) file."""
    if not OPENPYXL_AVAILABLE:
        return f"[XLSX file skipped - openpyxl not available: {file_path.name}]"
    try:
        workbook = load_workbook(str(file_path))
        text_parts = []
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_parts.append(f"--- Sheet: {sheet_name} ---")
            for row in sheet.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else "" for cell in row]
                text_parts.append(" | ".join(row_values))
        return "\n".join(text_parts)
    except Exception as e:
        return f"[Error extracting XLSX {file_path.name}: {str(e)}]"


@dataclass
class RetrievedChunk:
    doc_id: str
    title: str
    text: str
    score: float


class KnowledgeComplianceAgent:
    def __init__(
        self,
        model_name: str = "llama3:8b",
        embedding_model: str = "nomic-embed-text",
        db_dir: str = str(DEFAULT_DB_DIR),
        collection_name: str = DEFAULT_COLLECTION,
    ) -> None:
        self.client = Client()
        self.model = model_name
        self.embedding_model = embedding_model
        self.llm_options = {
            "temperature": 0.1,
            "top_p": 0.9,
            "num_predict": 350,
        }

        self.db_path = Path(db_dir)
        self.db_path.mkdir(parents=True, exist_ok=True)
        self.chroma = chromadb.PersistentClient(path=str(self.db_path))
        self.collection = self.chroma.get_or_create_collection(name=collection_name)

    def _embed_texts(self, texts: List[str]) -> List[List[float]]:
        result = self.client.embed(model=self.embedding_model, input=texts)
        return result.get("embeddings", [])

    def sync_knowledge_base(self, kb_dir: str, rebuild: bool = False) -> Dict[str, int]:
        """Sync folder files into the persistent vector DB.

        Supports: .txt, .md, .pdf (with OCR), .ppt, .pptx, .xlsx
        - New/updated files are upserted
        - For each synced file, old chunks are replaced
        - If rebuild=True, the whole collection is dropped/rebuilt
        """
        root = Path(kb_dir)
        root.mkdir(parents=True, exist_ok=True)

        if rebuild:
            self.chroma.delete_collection(self.collection.name)
            self.collection = self.chroma.get_or_create_collection(name=self.collection.name)

        files = [p for p in root.rglob("*") if p.is_file() and p.suffix.lower() in SUPPORTED_EXTENSIONS]
        files_synced = 0
        chunks_indexed = 0

        for file_path in files:
            # Extract text based on file type
            suffix = file_path.suffix.lower()
            if suffix == ".pdf":
                text = _extract_text_from_pdf(file_path)
            elif suffix == ".pptx":
                text = _extract_text_from_pptx(file_path)
            elif suffix == ".ppt":
                text = _extract_text_from_ppt(file_path)
            elif suffix == ".xlsx":
                text = _extract_text_from_xlsx(file_path)
            else:  # .txt, .md
                text = file_path.read_text(encoding="utf-8", errors="ignore")

            chunks = _chunk_text(text)
            if not chunks:
                continue

            source_path = str(file_path.resolve())

            # Replace old chunks for this file
            self.collection.delete(where={"source_path": source_path})

            ids = [f"{source_path}::chunk::{i}" for i in range(len(chunks))]
            titles = [f"{file_path.stem} (chunk {i + 1})" for i in range(len(chunks))]
            metadatas = [
                {
                    "source_path": source_path,
                    "title": titles[i],
                    "chunk_index": i,
                }
                for i in range(len(chunks))
            ]
            embeddings = self._embed_texts(chunks)

            self.collection.upsert(
                ids=ids,
                documents=chunks,
                metadatas=metadatas,
                embeddings=embeddings,
            )

            files_synced += 1
            chunks_indexed += len(chunks)

        return {"files_synced": files_synced, "chunks_indexed": chunks_indexed}

    def retrieve(self, query: str, top_k: int = 3) -> List[RetrievedChunk]:
        query_embedding = self._embed_texts([query])[0]
        result = self.collection.query(
            query_embeddings=[query_embedding],
            n_results=top_k,
            include=["documents", "metadatas", "distances"],
        )

        documents = (result.get("documents") or [[]])[0]
        metadatas = (result.get("metadatas") or [[]])[0]
        distances = (result.get("distances") or [[]])[0]

        chunks: List[RetrievedChunk] = []
        for i, doc in enumerate(documents):
            md = metadatas[i] if i < len(metadatas) else {}
            dist = distances[i] if i < len(distances) else 1.0
            score = 1.0 / (1.0 + float(dist))
            doc_id = f"R{i + 1}"
            title = md.get("title", "retrieved_chunk")
            chunks.append(RetrievedChunk(doc_id=doc_id, title=title, text=doc, score=score))

        return chunks

    def _build_grounded_prompt(self, user_input: str, intent: str, entities: List[str], chunks: List[RetrievedChunk]) -> str:
        context_lines = [f"[{c.doc_id}] {c.title}: {c.text}" for c in chunks]
        context = "\n".join(context_lines) if context_lines else "No context found."
        entities_line = ", ".join(entities) if entities else "none"

        return f"""You are Agent 2 (Knowledge & Compliance).

Task:
- Answer ONLY from the provided context.
- If context is insufficient, say exactly: 'I need more approved product information to answer safely.'
- Keep answer concise and factual.
- Include citation ids used in square brackets (example: [R1]).

User question: {user_input}
Intent from Agent 1: {intent}
Entities from Agent 1: {entities_line}

Context:
{context}

Return valid JSON only:
{{
  "answer": "<grounded response>",
  "citations": ["R1", "R2"]
}}"""

    def _call_llm(self, prompt: str) -> Dict[str, Any]:
        try:
            response = self.client.chat(
                model=self.model,
                messages=[{"role": "user", "content": prompt}],
                options=self.llm_options,
            )
            content = response["message"]["content"].strip()
            start = content.find("{")
            end = content.rfind("}")
            if start != -1 and end != -1 and end > start:
                content = content[start:end + 1]

            parsed = json.loads(content)
            citations = parsed.get("citations", [])
            if not isinstance(citations, list):
                citations = [str(citations)]
            return {
                "answer": parsed.get("answer", ""),
                "citations": [str(c) for c in citations],
            }
        except Exception:
            return {
                "answer": "I need more approved product information to answer safely.",
                "citations": [],
            }

    def _check_compliance(self, answer: str) -> List[str]:
        flags: List[str] = []
        lower = answer.lower()
        forbidden_patterns = [
            ("guaranteed", "overclaim_guarantee"),
            ("100% safe", "absolute_safety_claim"),
            ("no side effects", "absolute_safety_claim"),
            ("cure", "unverified_cure_claim"),
            ("off-label", "off_label_reference"),
        ]
        for phrase, flag in forbidden_patterns:
            if phrase in lower:
                flags.append(flag)
        return sorted(set(flags))

    def _compute_metrics(
        self,
        retrieved: List[RetrievedChunk],
        cited_ids: List[str],
        compliance_flags: List[str],
        answer: str,
    ) -> Dict[str, float]:
        retrieval_hit_quality = 0.0
        if retrieved:
            retrieval_hit_quality = sum(c.score for c in retrieved) / len(retrieved)

        retrieved_ids = {c.doc_id for c in retrieved}
        citation_coverage = len(retrieved_ids.intersection(set(cited_ids))) / len(retrieved_ids) if retrieved_ids else 0.0

        hallucination_rate = 1.0 if answer.strip() and not cited_ids else 0.0
        compliance_violation_rate = 1.0 if compliance_flags else 0.0

        return {
            "retrieval_hit_quality": round(retrieval_hit_quality, 4),
            "citation_coverage": round(citation_coverage, 4),
            "hallucination_rate": round(hallucination_rate, 4),
            "compliance_violation_rate": round(compliance_violation_rate, 4),
        }

    def process(self, user_input: str, intent: str, entities: List[str] | None = None, top_k: int = 3) -> Dict[str, Any]:
        entities = entities or []
        retrieved = self.retrieve(user_input, top_k=top_k)
        prompt = self._build_grounded_prompt(user_input, intent, entities, retrieved)
        llm_result = self._call_llm(prompt)

        answer = llm_result.get("answer", "")
        citations = [c for c in llm_result.get("citations", []) if isinstance(c, str)]
        compliance_flags = self._check_compliance(answer)
        metrics = self._compute_metrics(retrieved, citations, compliance_flags, answer)

        top_score = retrieved[0].score if retrieved else 0.0
        confidence = max(0.0, min(1.0, top_score - 0.2 * len(compliance_flags)))

        return {
            "answer": answer,
            "citations": citations,
            "confidence": round(confidence, 4),
            "compliance_flags": compliance_flags,
            "metrics": metrics,
            "next_agent": "conversation_manager",
            "retrieved_chunks": [
                {
                    "doc_id": c.doc_id,
                    "title": c.title,
                    "score": round(c.score, 4),
                }
                for c in retrieved
            ],
        }


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Agent 2 - Knowledge & Compliance (RAG + Vector DB)")
    parser.add_argument("--question", default="What is the recommended dosage for Cardiomax in adults?")
    parser.add_argument("--intent", default="dosage_question")
    parser.add_argument("--entities", nargs="*", default=["Cardiomax"])
    parser.add_argument("--kb-dir", default=str(DEFAULT_KB_DIR), help="Folder with .txt/.md files")
    parser.add_argument("--db-dir", default=str(DEFAULT_DB_DIR), help="Persistent Chroma DB folder")
    parser.add_argument("--collection", default=DEFAULT_COLLECTION, help="Chroma collection name")
    parser.add_argument("--embedding-model", default="nomic-embed-text", help="Ollama embedding model")
    parser.add_argument("--rebuild", action="store_true", help="Rebuild vector collection from scratch")
    args = parser.parse_args()

    agent = KnowledgeComplianceAgent(
        embedding_model=args.embedding_model,
        db_dir=args.db_dir,
        collection_name=args.collection,
    )

    sync_info = agent.sync_knowledge_base(args.kb_dir, rebuild=args.rebuild)
    print(json.dumps({"sync": sync_info}, indent=2))

    result = agent.process(
        user_input=args.question,
        intent=args.intent,
        entities=args.entities,
    )
    print(json.dumps(result, indent=2))
